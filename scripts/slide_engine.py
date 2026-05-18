"""
slide_engine.py — shared helpers for all phase PPTX builders.
Academic M.Tech style: dark navy header bar, white body, clean monospace for code/math.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree
import copy, io

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0A, 0x23, 0x4F)   # header bar / section accent
BLUE      = RGBColor(0x16, 0x4E, 0x9C)   # secondary accent
TEAL      = RGBColor(0x00, 0x7A, 0x87)   # highlight boxes
ORANGE    = RGBColor(0xD4, 0x5D, 0x00)   # warning / change marker
GREEN     = RGBColor(0x17, 0x6B, 0x38)   # positive / result
RED       = RGBColor(0xB5, 0x19, 0x19)   # negative / old
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF4, 0xF6, 0xFA)
NEAR_BLK  = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_BLU = RGBColor(0xD8, 0xE8, 0xF8)   # content box tint
LIGHT_GRN = RGBColor(0xD4, 0xED, 0xDA)
LIGHT_ORG = RGBColor(0xFD, 0xE8, 0xD0)
MID_GRAY  = RGBColor(0x5A, 0x6A, 0x7A)
LINE_GRAY = RGBColor(0xCC, 0xD5, 0xDF)

# ── Slide setup ───────────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)
MARGIN_L = Inches(0.45)
CONTENT_W = Inches(12.43)

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

# ── Background ────────────────────────────────────────────────────────────────
def set_bg(slide, color=OFF_WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

# ── Shapes ────────────────────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=None, line_color=None, line_w=Pt(0.75), radius=0):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def divider(slide, t, color=LINE_GRAY, thickness=Pt(0.5)):
    ln = slide.shapes.add_shape(1, MARGIN_L, t, CONTENT_W, thickness)
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    return ln

# ── Text helpers ──────────────────────────────────────────────────────────────
def _run(para, text, size, bold=False, italic=False, color=NEAR_BLK, mono=False, underline=False):
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.underline = underline
    run.font.color.rgb = color
    if mono:
        run.font.name = "Courier New"
    return run

def textbox(slide, text, l, t, w, h,
            size=13, bold=False, italic=False, color=NEAR_BLK,
            align=PP_ALIGN.LEFT, wrap=True, mono=False,
            line_spacing=None):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    lines = text.split('\n')
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        _run(p, line, size, bold=bold, italic=italic, color=color, mono=mono)
    return txb

def rich_textbox(slide, segments, l, t, w, h, wrap=True, line_spacing=None):
    """
    segments: list of (text, size, bold, italic, color, mono, new_line)
    new_line=True starts a new paragraph before this segment.
    """
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    for (text, size, bold, italic, color, mono, new_para) in segments:
        if new_para and para.runs:
            para = tf.add_paragraph()
        if line_spacing and new_para:
            para.line_spacing = line_spacing
        _run(para, text, size, bold=bold, italic=italic, color=color, mono=mono)
    return txb

# ── Standard slide chrome ─────────────────────────────────────────────────────
def slide_chrome(slide, phase_label, slide_title, subtitle=None):
    """
    Top bar: dark navy.  Phase tag left.  Title center-left.
    """
    set_bg(slide, OFF_WHITE)
    # top bar
    rect(slide, 0, 0, W, Inches(0.72), fill=NAVY)
    # phase pill
    rect(slide, Inches(0.35), Inches(0.1), Inches(1.65), Inches(0.52), fill=BLUE)
    textbox(slide, phase_label,
            Inches(0.35), Inches(0.1), Inches(1.65), Inches(0.52),
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # slide title
    textbox(slide, slide_title,
            Inches(2.2), Inches(0.08), Inches(10.8), Inches(0.58),
            size=20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        textbox(slide, subtitle,
                Inches(2.2), Inches(0.6), Inches(10.8), Inches(0.28),
                size=11, italic=True, color=RGBColor(0xB0, 0xC8, 0xE8))
    # bottom rule
    rect(slide, 0, H - Inches(0.28), W, Inches(0.28), fill=NAVY)
    return slide

def footer(slide, text, right_text=None):
    textbox(slide, text,
            MARGIN_L, H - Inches(0.25), Inches(7), Inches(0.22),
            size=8, color=RGBColor(0xC0, 0xCE, 0xDE))
    if right_text:
        textbox(slide, right_text,
                Inches(6.5), H - Inches(0.25), Inches(6.5), Inches(0.22),
                size=8, color=RGBColor(0xC0, 0xCE, 0xDE), align=PP_ALIGN.RIGHT)

def section_label(slide, text, l, t, w, h, color=NAVY):
    rect(slide, l, t, Inches(0.06), h, fill=color)
    textbox(slide, text, l + Inches(0.14), t, w - Inches(0.14), h,
            size=11, bold=True, color=color)

# ── Callout / key-finding box ─────────────────────────────────────────────────
def callout_box(slide, title, body, l, t, w, h,
                accent=TEAL, bg=LIGHT_BLU, title_size=12, body_size=11.5):
    rect(slide, l, t, w, h, fill=bg, line_color=accent, line_w=Pt(1.5))
    rect(slide, l, t, Inches(0.07), h, fill=accent)
    textbox(slide, title, l + Inches(0.17), t + Inches(0.08),
            w - Inches(0.22), Inches(0.32),
            size=title_size, bold=True, color=accent)
    textbox(slide, body, l + Inches(0.17), t + Inches(0.35),
            w - Inches(0.22), h - Inches(0.42),
            size=body_size, color=NEAR_BLK)

def result_box(slide, label, value, l, t, w, h, accent=GREEN):
    rect(slide, l, t, w, h, fill=LIGHT_GRN if accent == GREEN else LIGHT_ORG,
         line_color=accent, line_w=Pt(1.2))
    textbox(slide, label, l + Inches(0.1), t + Inches(0.06),
            w - Inches(0.15), Inches(0.28),
            size=10, bold=True, color=accent)
    textbox(slide, value, l + Inches(0.1), t + Inches(0.3),
            w - Inches(0.15), h - Inches(0.38),
            size=13, bold=True, color=NEAR_BLK, align=PP_ALIGN.CENTER)

# ── Image helper ─────────────────────────────────────────────────────────────
def add_img(slide, path, l, t, w, h):
    from pathlib import Path
    p = Path(path)
    if p.exists():
        slide.shapes.add_picture(str(p), l, t, w, h)
    else:
        r = rect(slide, l, t, w, h, fill=RGBColor(0xE0, 0xE8, 0xF0),
                 line_color=LINE_GRAY)
        textbox(slide, f"[{p.name}]", l, t + h//2 - Inches(0.2), w, Inches(0.4),
                size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ── Table helper ──────────────────────────────────────────────────────────────
def add_table(slide, headers, rows, l, t, w, h,
              header_fill=NAVY, header_color=WHITE, header_size=11,
              row_size=11, alt_fill=RGBColor(0xEA, 0xF0, 0xF8)):
    cols = len(headers)
    total_rows = 1 + len(rows)
    tbl = slide.shapes.add_table(total_rows, cols, l, t, w, h).table
    col_w = w // cols
    for i in range(cols):
        tbl.columns[i].width = col_w

    # header row
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        tf = cell.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = hdr
        run.font.size = Pt(header_size)
        run.font.bold = True
        run.font.color.rgb = header_color

    # data rows
    for ri, row in enumerate(rows):
        fill_color = alt_fill if ri % 2 == 1 else WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = fill_color
            tf = cell.text_frame
            align = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            tf.paragraphs[0].alignment = align
            run = tf.paragraphs[0].add_run()
            run.text = str(val)
            run.font.size = Pt(row_size)
            run.font.color.rgb = NEAR_BLK
    return tbl
